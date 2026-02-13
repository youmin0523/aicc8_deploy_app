from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define some colors
    DARK_BG = RGBColor(18, 18, 18)
    TEXT_WHITE = RGBColor(255, 255, 255)
    ACCENT_BLUE = RGBColor(74, 144, 226)
    ACCENT_EMERALD = RGBColor(80, 200, 120)

    # Helper to add a slide with title and content
    def add_slide(title_text, content_text, layout_index=1):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Dark Background for all slides
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
        title.text_frame.paragraphs[0].font.bold = True
        
        # Content
        if content_text:
            content = slide.placeholders[1]
            content.text = content_text
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.color.rgb = TEXT_WHITE
                paragraph.font.size = Pt(18)
        
        return slide

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    title = slide.shapes.title
    title.text = "AICC8 Deploy App: V2 Upgrade Project"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Next-Generation Task Management System\nProject Report & Technical Overview"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.color.rgb = TEXT_WHITE

    # Slide 2: Project Overview
    add_slide(
        "1. Project Overview",
        "• Mission: Upgrade existing V1 Task Manager to a Premium V2 Dashboard.\n"
        "• Key Goal: Gamification, Enhanced UX, and Data Integrity.\n"
        "• Core Concept: 'Space Mission' theme with immersive animations.\n"
        "• Architecture: Hybrid System supporting both V1 and V2 data structures simultaneously."
    )

    # Slide 3: Key Features (V2)
    add_slide(
        "2. Key Features (V2 Dashboard)",
        "DASHBOARD & UX\n"
        "• Dark Mode & Glassmorphism Design System.\n"
        "• Intro: Cinematic Space Launch Animation (Zoom-in effect).\n"
        "• Version Loader: Seamless transition between V1 and V2.\n\n"
        "TASK MANAGEMENT\n"
        "• Neon Orb System: Visual status indicators (Today/Tomorrow).\n"
        "• Audit History: Detailed log of value changes (Title, Desc, Status).\n"
        "• Integrated Calendar: Sunday-start layout with drag-free navigation."
    )

    # Slide 4: Technical Stack
    add_slide(
        "3. Technical Stack",
        "FRONTEND\n"
        "• Framework: React 19 + Vite\n"
        "• State Management: Redux Toolkit (Slices for Auth, Tasks, Modals)\n"
        "• Styling: TailwindCSS + Vanilla CSS (No Styled-Components)\n"
        "• Libraries: react-calendar, react-toastify, react-icons, react-oauth/google\n\n"
        "BACKEND\n"
        "• Runtime: Node.js + Express\n"
        "• Database: PostgreSQL (Hybrid Schema: tasks & tasks_v2)\n"
        "• Security: Google OAuth 2.0 Integration"
    )

    # Slide 5: System Architecture
    add_slide(
        "4. System Architecture",
        "HYBRID DATA MODEL\n"
        "• V1 Compatibility: Legacy 'tasks' table retained.\n"
        "• V2 Innovation: New 'tasks_v2' table with 'categoryid', 'updated_at'.\n"
        "• Unified View: 'UNION ALL' query to display V1 & V2 tasks in one list.\n\n"
        "DATA FLOW\n"
        "• User Action -> Redux Thunk (Async) -> API Controller -> PostgreSQL.\n"
        "• Audit Log: 'logTaskChange' function tracks every modification."
    )

    # Slide 6: Troubleshooting & Improvements
    add_slide(
        "5. Major Troubleshooting History",
        "• Rendering: Fixed 'react-calendar' v6 blank screen issue by switching to 'gregory' type.\n"
        "• API: Solved 'Empty Body' issue in PATCH requests by fixing JSON stringification.\n"
        "• Stability: Implemented Fail-Safe logic for Audit Logging (Log failure doesn't block Task update).\n"
        "• UX: Resolved 'Black Screen' on version transition with persistent Overlay Loader."
    )

    # Slide 7: Future Roadmap
    add_slide(
        "6. Future Roadmap",
        "• Mobile Optimization: Fully responsive layout for mobile devices.\n"
        "• Advanced Analytics: Productivity charts based on History data.\n"
        "• Team Collaboration: Shared workspaces and task delegation.\n"
        "• AI Integration: Smart task suggestions and auto-categorization."
    )

    # Save
    prs.save('AICC8_Project_Report.pptx')
    print("Presentation saved successfully as AICC8_Project_Report.pptx")

if __name__ == "__main__":
    create_presentation()
